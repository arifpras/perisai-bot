"""Generate multi-page PDF and PowerPoint presentations from LLM responses.

Supports:
- PDF generation with reportlab
- PowerPoint (PPTX) generation with python-pptx
- Automatic slide/page parsing from markdown-style content
- Professional formatting and styling
- Image embedding
- Table generation

Usage:
    from presentation_generator import generate_pdf, generate_pptx
    
    pdf_bytes = generate_pdf(content, title="My Presentation")
    pptx_bytes = generate_pptx(content, title="My Presentation")
"""

import io
import logging
import re
from typing import List, Dict, Optional, Tuple
from datetime import datetime

logger = logging.getLogger(__name__)

# Optional imports - graceful degradation if not installed
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor, white, black, grey
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY, TA_RIGHT
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False
    logger.warning("reportlab not installed - PDF generation disabled")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed - PowerPoint generation disabled")


# Styling constants
COLORS = {
    'primary': '#1F4788',      # Dark blue
    'secondary': '#2E5C8A',    # Medium blue
    'accent': '#E74C3C',       # Red
    'text': '#2C3E50',         # Dark gray
    'light_gray': '#ECF0F1',   # Light gray
    'white': '#FFFFFF',
}

FONTS = {
    'title': 24,
    'heading': 16,
    'subheading': 14,
    'body': 11,
    'small': 10,
}


class SlideContent:
    """Represents a single slide/page of content."""
    
    def __init__(self, title: str, content: List[str], slide_type: str = "content"):
        """
        Args:
            title: Slide title
            content: List of content lines/paragraphs
            slide_type: 'title', 'content', 'bullets', 'table', 'conclusion'
        """
        self.title = title
        self.content = content
        self.slide_type = slide_type


def parse_presentation_content(text: str) -> List[SlideContent]:
    """Parse LLM response into individual slides.
    
    Expected format:
    # Main Title
    
    ## Slide 1: Title
    Content paragraph 1
    Content paragraph 2
    
    - Bullet point 1
    - Bullet point 2
    
    ## Slide 2: Another Title
    More content...
    
    Args:
        text: Raw text content from LLM
        
    Returns:
        List of SlideContent objects
    """
    slides = []
    
    # Split by ## (slide headers)
    slide_pattern = r'##\s+(.+?)(?=##\s+|\Z)'
    matches = re.finditer(slide_pattern, text, re.DOTALL)
    
    for match in matches:
        slide_text = match.group(1).strip()
        lines = slide_text.split('\n')
        
        if not lines:
            continue
        
        # First line is title
        title = lines[0].strip().rstrip(':')
        
        # Rest is content
        content_lines = [line.strip() for line in lines[1:] if line.strip()]
        
        # Detect slide type
        slide_type = "content"
        if any(line.startswith('-') or line.startswith('*') for line in content_lines):
            slide_type = "bullets"
        elif any('|' in line for line in content_lines):
            slide_type = "table"
        elif re.search(r'conclusion|summary|key point', title.lower()):
            slide_type = "conclusion"
        
        slides.append(SlideContent(title, content_lines, slide_type))
    
    return slides


def extract_title(text: str) -> str:
    """Extract main title from content (first # heading)."""
    match = re.search(r'^#\s+(.+?)$', text, re.MULTILINE)
    return match.group(1) if match else "Presentation"


def extract_bullets(content_lines: List[str]) -> List[str]:
    """Extract bullet points from content lines."""
    bullets = []
    for line in content_lines:
        if re.match(r'^\s*[-*•]\s+', line):
            bullet = re.sub(r'^\s*[-*•]\s+', '', line)
            bullets.append(bullet)
    return bullets


def generate_pdf(
    content: str,
    title: Optional[str] = None,
    author: str = "PerisAI Bot",
    output_format: str = "letter"
) -> Optional[bytes]:
    """Generate PDF presentation from content.
    
    Args:
        content: Content to convert to PDF
        title: Presentation title (auto-extracted if not provided)
        author: Document author name
        output_format: 'letter' or 'a4'
        
    Returns:
        PDF file as bytes, or None if reportlab unavailable
    """
    if not HAS_REPORTLAB:
        logger.error("reportlab not installed - PDF generation unavailable")
        return None
    
    try:
        # Parse content into slides
        if not title:
            title = extract_title(content)
        
        slides = parse_presentation_content(content)
        
        # Create PDF document
        pdf_buffer = io.BytesIO()
        page_size = A4 if output_format == "a4" else letter
        doc = SimpleDocTemplate(
            pdf_buffer,
            pagesize=page_size,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=1*inch,
            bottomMargin=0.75*inch,
            title=title,
            author=author,
        )
        
        # Build story
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=FONTS['title'],
            textColor=HexColor(COLORS['primary']),
            spaceAfter=12,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        slide_heading_style = ParagraphStyle(
            'SlideHeading',
            parent=styles['Heading2'],
            fontSize=FONTS['heading'],
            textColor=HexColor(COLORS['secondary']),
            spaceAfter=12,
            spaceBefore=6,
            fontName='Helvetica-Bold',
            borderPadding=6,
            borderColor=HexColor(COLORS['secondary']),
            borderWidth=1,
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=FONTS['body'],
            alignment=TA_JUSTIFY,
            spaceAfter=10,
        )
        
        # Add title page
        story.append(Spacer(1, 1*inch))
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(
            f"Generated by PerisAI Bot on {datetime.now().strftime('%B %d, %Y')}",
            styles['Normal']
        ))
        story.append(PageBreak())
        
        # Add slides
        for i, slide in enumerate(slides):
            # Slide heading
            story.append(Paragraph(slide.title, slide_heading_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Content based on slide type
            if slide.slide_type == "bullets":
                bullets = extract_bullets(slide.content)
                for bullet in bullets:
                    bullet_text = f"• {bullet}"
                    story.append(Paragraph(bullet_text, body_style))
                    story.append(Spacer(1, 0.1*inch))
            
            elif slide.slide_type == "table":
                # Try to parse as table
                table_rows = []
                for line in slide.content:
                    if '|' in line:
                        cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                        if cells:
                            table_rows.append(cells)
                
                if table_rows:
                    table = Table(table_rows)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor(COLORS['secondary'])),
                        ('TEXTCOLOR', (0, 0), (-1, 0), white),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), FONTS['body']),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), HexColor(COLORS['light_gray'])),
                        ('GRID', (0, 0), (-1, -1), 1, HexColor(COLORS['text'])),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 0.2*inch))
            
            else:
                # Regular content
                for line in slide.content:
                    if line:
                        story.append(Paragraph(line, body_style))
                        story.append(Spacer(1, 0.1*inch))
            
            # Page break between slides
            if i < len(slides) - 1:
                story.append(PageBreak())
        
        # Build PDF
        doc.build(story)
        pdf_buffer.seek(0)
        return pdf_buffer.read()
    
    except Exception as e:
        logger.error(f"PDF generation failed: {e}", exc_info=True)
        return None


def generate_pptx(
    content: str,
    title: Optional[str] = None,
    author: str = "PerisAI Bot"
) -> Optional[bytes]:
    """Generate PowerPoint presentation from content.
    
    Args:
        content: Content to convert to PPTX
        title: Presentation title (auto-extracted if not provided)
        author: Document author name
        
    Returns:
        PPTX file as bytes, or None if python-pptx unavailable
    """
    if not HAS_PPTX:
        logger.error("python-pptx not installed - PowerPoint generation unavailable")
        return None
    
    try:
        # Parse content into slides
        if not title:
            title = extract_title(content)
        
        slides = parse_presentation_content(content)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Set document properties
        prs.core_properties.title = title
        prs.core_properties.author = author
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"Generated by {author}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Format title
        title_shape.text_frame.paragraphs[0].font.size = Pt(FONTS['title'])
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)  # Primary color
        
        # Add content slides
        for slide_num, slide_content in enumerate(slides):
            # Choose layout based on content type
            if slide_content.slide_type == "bullets":
                layout = prs.slide_layouts[1]  # Title and Content
            elif slide_content.slide_type == "table":
                layout = prs.slide_layouts[5]  # Blank layout (will add table manually)
            else:
                layout = prs.slide_layouts[1]  # Title and Content
            
            slide = prs.slides.add_slide(layout)
            
            # Add title
            if slide.shapes.title:
                slide.shapes.title.text = slide_content.title
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(FONTS['heading'])
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 92, 138)  # Secondary color
            
            # Add content
            if slide_content.slide_type == "bullets":
                # Use content placeholder for bullets
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    
                    bullets = extract_bullets(slide_content.content)
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(FONTS['body'])
            
            elif slide_content.slide_type == "table":
                # Parse table from content
                table_rows = []
                for line in slide_content.content:
                    if '|' in line:
                        cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                        if cells:
                            table_rows.append(cells)
                
                if table_rows:
                    # Add table shape
                    rows = len(table_rows)
                    cols = len(table_rows[0])
                    
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    
                    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Fill table with content
                    for row_idx, row_data in enumerate(table_rows):
                        for col_idx, cell_data in enumerate(row_data):
                            cell = table_shape.cell(row_idx, col_idx)
                            cell.text = str(cell_data)
                            
                            # Format header row
                            if row_idx == 0:
                                cell.text_frame.paragraphs[0].font.bold = True
                                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            else:
                # Regular content
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    
                    for i, line in enumerate(slide_content.content):
                        if line:
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            p.text = line
                            p.level = 0
                            p.font.size = Pt(FONTS['body'])
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.read()
    
    except Exception as e:
        logger.error(f"PowerPoint generation failed: {e}", exc_info=True)
        return None


def get_file_extension_and_bytes(
    format_type: str,
    content: str,
    title: Optional[str] = None
) -> Tuple[Optional[str], Optional[bytes]]:
    """Generate presentation in specified format.
    
    Args:
        format_type: 'pdf' or 'pptx'
        content: Content to convert
        title: Presentation title
        
    Returns:
        Tuple of (file_extension, file_bytes)
    """
    if format_type.lower() == "pdf":
        pdf_bytes = generate_pdf(content, title)
        return (".pdf", pdf_bytes) if pdf_bytes else (None, None)
    
    elif format_type.lower() in ["pptx", "powerpoint"]:
        pptx_bytes = generate_pptx(content, title)
        return (".pptx", pptx_bytes) if pptx_bytes else (None, None)
    
    else:
        logger.error(f"Unknown format type: {format_type}")
        return (None, None)


# Integration example for telegram_bot.py:
"""
In ask_kei() or ask_kin() functions, check for presentation requests:

    question_lower = question.lower()
    
    # Check for PDF generation request
    if any(keyword in question_lower for keyword in ['generate pdf', 'pdf presentation', 'create pdf']):
        try:
            answer = await ask_kei(question)  # Get content from Kei
            ext, file_bytes = get_file_extension_and_bytes('pdf', answer, title=question)
            if file_bytes:
                # Send PDF file to user
                await update.message.reply_document(
                    document=io.BytesIO(file_bytes),
                    filename=f"presentation{ext}",
                    caption="📄 PDF Presentation"
                )
                return
        except Exception as e:
            logger.error(f"PDF generation failed: {e}")
    
    # Check for PowerPoint generation request
    elif any(keyword in question_lower for keyword in ['generate pptx', 'create powerpoint', 'powerpoint presentation']):
        try:
            answer = await ask_kei(question)  # Get content from Kei
            ext, file_bytes = get_file_extension_and_bytes('pptx', answer, title=question)
            if file_bytes:
                # Send PPTX file to user
                await update.message.reply_document(
                    document=io.BytesIO(file_bytes),
                    filename=f"presentation{ext}",
                    caption="📊 PowerPoint Presentation"
                )
                return
        except Exception as e:
            logger.error(f"PowerPoint generation failed: {e}")
"""
